######-----Import Dash-----#####
import dash
from dash import dcc
from dash import html, callback
import dash_bootstrap_components as dbc
from dash.dependencies import Input, Output
from dash.exceptions import PreventUpdate

from pptx import Presentation
from pptx.util import Cm, Pt

import os

import pandas as pd
import numpy as np
import plotly.graph_objects as go
import plotly.io as pio

from datetime import date
from dateutil.relativedelta import relativedelta

dash.register_page(__name__, name='Download')

##-----page

## -----LAYOUT-----
layout = html.Div([
                html.Br(),
                
                html.Div([
                    html.Button('Download Slides', 
                                      id='download_button', 
                                      n_clicks=0, 
                                      style={'fontSize': 15, 'color':'#2a3f5f','font-family':'Verdana','display': 'inline-block'}), 
                    dcc.Download(id='download'),
                    ]),

                html.Br(),
                html.P(id='download_date', style={'fontSize': 15, 'color':'#2a3f5f','font-family':'Verdana'}),

    html.Footer('ABL',
            style={'textAlign': 'center', 
                   'fontSize': 20, 
                   'background-color':'#2a3f5f',
                   'color':'white',
                   'position': 'fixed',
                    'bottom': '0',
                    'width': '100%'})

    ], style={
        'paddingLeft':'10px',
        'paddingRight':'10px',
    })

#### Callback Auto Update Chart & Data

@callback(
    [Output('download', 'data'),
     Output('download_date', 'children')],
    [Input('download_button', 'n_clicks'),
     Input('store', 'data')]
)
def update_charts(n_clicks, data):
    if n_clicks == 0:
        raise PreventUpdate
    ######################
    # Pre-Processing
    ######################
    df_mv = pd.DataFrame(data['Mother Vessel'])
    df_mv = df_mv.replace(r'^\s*$', np.nan, regex=True)
    date_mv = ['Keel Laying', 'Year of Build', 
                'Annual Survey Last Date', 'Annual Survey Due Date',
                'Special Survey Last Date', 'Special Survey Due Date',
                'Intermediate Survey Last Date', 'Intermediate Survey Due Date',
                'Docking Survey Last Date', 'Docking Survey Due Date',
                'Tail Shaft Survey Last Date', 'Tail Shaft Survey Due Date',
                'Boiler Survey Last Date', 'Boiler Survey Due Date',
                'COC Finding Date', 'COC Due Date'] 
    df_mv[date_mv] = df_mv[date_mv].apply(pd.to_datetime, format='%d/%m/%Y')

    df_cts = pd.DataFrame(data['Floating Crane'])
    df_cts = df_cts.replace(r'^\s*$', np.nan, regex=True)
    date_cts = ['Keel Laying', 'Year of Build', 
            'Annual Survey Last Date', 'Annual Survey Due Date', 
            'Special Survey Last Date', 'Special Survey Due Date',
            'Intermediate Survey Last Date', 'Intermediate Survey Due Date',
            'Docking Survey Last Date', 'Docking Survey Due Date',
            'Cargo Gear Survey Last Date', 'Cargo Gear Survey Due Date',
            'COC Finding Date', 'COC Due Date'] 
    df_cts[date_cts] = df_cts[date_cts].apply(pd.to_datetime, format='%d/%m/%Y')

    df_tb = pd.DataFrame(data['Tugboat'])
    df_tb = df_tb.replace(r'^\s*$', np.nan, regex=True)
    date_tb = ['Keel Laying', 'Year of Build', 
            'Annual Survey Last Date', 'Annual Survey Due Date',
            'Special Survey Last Date', 'Special Survey Due Date',
            'Intermediate Survey Last Date', 'Intermediate Survey Due Date',
            'Docking Survey Last Date', 'Docking Survey Due Date',
            'Tail Shaft Survey Last Date', 'Tail Shaft Survey Due Date',
            'COC Finding Date', 'COC Due Date'] 
    df_tb[date_tb] = df_tb[date_tb].apply(pd.to_datetime, format='%d/%m/%Y')

    df_ba = pd.DataFrame(data['Barge'])
    df_ba = df_ba.replace(r'^\s*$', np.nan, regex=True)
    date_ba = ['Keel Laying', 
            'Annual Survey Last Date', 'Annual Survey Due Date',
            'Special Survey Last Date', 'Special Survey Due Date',
            'Intermediate Survey Last Date', 'Intermediate Survey Due Date',
            'Docking Survey Last Date', 'Docking Survey Due Date',
            'COC Finding Date', 'COC Due Date'] 
    df_ba[date_ba] = df_ba[date_ba].apply(pd.to_datetime, format='%d/%m/%Y')

    # New dataframe for remaining days
    diff_mv = ['Annual Survey Due Date', 'Special Survey Due Date', 'Intermediate Survey Due Date', 'Docking Survey Due Date', 'Tail Shaft Survey Due Date', 'Boiler Survey Due Date', 'COC Due Date']
    df_diff_mv = df_mv[['Mother Vessel Name', 'Class', 'Remarks']+diff_mv].copy()

    diff_cts = ['Annual Survey Due Date', 'Special Survey Due Date', 'Intermediate Survey Due Date', 'Docking Survey Due Date', 'Cargo Gear Survey Due Date', 'COC Due Date']
    df_diff_cts = df_cts[['CTS Name', 'Class', 'Remarks']+diff_cts].copy()

    diff_tb = ['Annual Survey Due Date', 'Special Survey Due Date', 'Intermediate Survey Due Date', 'Docking Survey Due Date', 'Tail Shaft Survey Due Date', 'COC Due Date']
    df_diff_tb = df_tb[['Tugboat Name', 'Class', 'Remarks']+diff_tb].copy()

    diff_ba = ['Annual Survey Due Date', 'Special Survey Due Date', 'Intermediate Survey Due Date', 'Docking Survey Due Date', 'COC Due Date']
    df_diff_ba = df_ba[['Barge Name', 'Class', 'Remarks']+diff_ba].copy()

    # Add new variable for remaining days
    diff_mv_name = ['Annual Survey Remain Days', 'Special Survey Remain Days', 'Intermediate Survey Remain Days', 'Docking Survey Remain Days', 'Tail Shaft Survey Remain Days', 'Boiler Survey Remain Days', 'COC Remain Days']
    diff_cts_name = ['Annual Survey Remain Days', 'Special Survey Remain Days', 'Intermediate Survey Remain Days', 'Docking Survey Remain Days', 'Cargo Gear Survey Remain Days', 'COC Remain Days']
    diff_tb_name = ['Annual Survey Remain Days', 'Special Survey Remain Days', 'Intermediate Survey Remain Days', 'Docking Survey Remain Days', 'Tail Shaft Survey Remain Days', 'COC Remain Days']
    diff_ba_name = ['Annual Survey Remain Days', 'Special Survey Remain Days', 'Intermediate Survey Remain Days', 'Docking Survey Remain Days', 'COC Remain Days']

    for (i,j) in zip(diff_mv, diff_mv_name):
        df_diff_mv[j] = (df_diff_mv[i]-pd.to_datetime('today')).dt.days

    for (i,j) in zip(diff_cts, diff_cts_name):
        df_diff_cts[j] = (df_diff_cts[i]-pd.to_datetime('today')).dt.days

    for (i,j) in zip(diff_tb, diff_tb_name):
        df_diff_tb[j] = (df_diff_tb[i]-pd.to_datetime('today')).dt.days

    for (i,j) in zip(diff_ba, diff_ba_name):
        df_diff_ba[j] = (df_diff_ba[i]-pd.to_datetime('today')).dt.days

    # Add new variable for remaining days (Months, Days)
    diff_mv_name2 = ['Annual Survey Remain Days2', 'Special Survey Remain Days2', 'Intermediate Survey Remain Days2', 'Docking Survey Remain Days2', 'Tail Shaft Survey Remain Days2', 'Boiler Survey Remain Days2', 'COC Remain Days2']
    diff_cts_name2 = ['Annual Survey Remain Days2', 'Special Survey Remain Days2', 'Intermediate Survey Remain Days2', 'Docking Survey Remain Days2', 'Cargo Gear Survey Remain Days2', 'COC Remain Days2']
    diff_tb_name2 = ['Annual Survey Remain Days2', 'Special Survey Remain Days2', 'Intermediate Survey Remain Days2', 'Docking Survey Remain Days2', 'Tail Shaft Survey Remain Days2', 'COC Remain Days2']
    diff_ba_name2 = ['Annual Survey Remain Days2', 'Special Survey Remain Days2', 'Intermediate Survey Remain Days2', 'Docking Survey Remain Days2', 'COC Remain Days2']

    # Function for calculating remaining days
    def date_diff(var):
        remaining =[]
        for i in var:
            if i is pd.NaT:
                remaining.append(pd.Timestamp('NaT'))
            else :
                delta = relativedelta(i, pd.to_datetime('today'))
                res_months = delta.months + (delta.years * 12)
                remaining.append((str(res_months)+ ' months, '+ str(delta.days)+ ' days'))
        return(remaining)

    for i,j in zip(diff_mv, diff_mv_name2):
        df_diff_mv[j] = date_diff(df_diff_mv[i])

    for i,j in zip(diff_cts, diff_cts_name2):
        df_diff_cts[j] = date_diff(df_diff_cts[i])

    for i,j in zip(diff_tb, diff_tb_name2):
        df_diff_tb[j] = date_diff(df_diff_tb[i])

    for i,j in zip(diff_ba, diff_ba_name2):
        df_diff_ba[j] = date_diff(df_diff_ba[i])

    # Rename vessel name column
    df_diff_mv = df_diff_mv.rename(columns={'Mother Vessel Name':'Vessel Name'})
    df_diff_cts = df_diff_cts.rename(columns={'CTS Name':'Vessel Name'})
    df_diff_tb = df_diff_tb.rename(columns={'Tugboat Name':'Vessel Name'})
    df_diff_ba = df_diff_ba.rename(columns={'Barge Name':'Vessel Name'})

    # Add emoji symbol for alerts
    df_diff_mv['Annual'] = df_diff_mv['Annual Survey Remain Days'].apply(lambda x:
                                                                        '❌' if x < 0 else (
                                                                            '⚠️' if x < 92 else '✔️'))
    old_mv = ['Special Survey Remain Days', 'Intermediate Survey Remain Days', 'Docking Survey Remain Days', 'Tail Shaft Survey Remain Days', 'Boiler Survey Remain Days', 'COC Remain Days']
    new_mv = ['Special', 'Intermediate', 'Docking', 'Tail Shaft', 'Boiler', 'COC']
    for i,j in zip(new_mv, old_mv):
        df_diff_mv[i] = df_diff_mv[j].apply(lambda x:
                                            '❌' if x < 0 else (
                                            '⚠️' if x < 183 else '✔️'))

    df_diff_cts['Annual'] = df_diff_cts['Annual Survey Remain Days'].apply(lambda x:
                                                                        '❌' if x < 0 else (
                                                                            '⚠️' if x < 92 else '✔️'))
    old_cts = ['Special Survey Remain Days', 'Intermediate Survey Remain Days', 'Docking Survey Remain Days', 'COC Remain Days']
    new_cts = ['Special', 'Intermediate', 'Docking', 'COC']
    for i,j in zip(new_cts, old_cts):
        df_diff_cts[i] = df_diff_cts[j].apply(lambda x:
                                            '❌' if x < 0 else (
                                            '⚠️' if x < 183 else '✔️'))

    df_diff_tb['Annual'] = df_diff_tb['Annual Survey Remain Days'].apply(lambda x:
                                                                        '❌' if x < 0 else (
                                                                            '⚠️' if x < 92 else '✔️'))
    old_tb = ['Special Survey Remain Days', 'Intermediate Survey Remain Days', 'Docking Survey Remain Days', 'Tail Shaft Survey Remain Days', 'COC Remain Days']
    new_tb = ['Special', 'Intermediate', 'Docking', 'Tail Shaft', 'COC']
    for i,j in zip(new_tb, old_tb):
        df_diff_tb[i] = df_diff_tb[j].apply(lambda x:
                                            '❌' if x < 0 else (
                                            '⚠️' if x < 183 else '✔️'))

    df_diff_ba['Annual'] = df_diff_ba['Annual Survey Remain Days'].apply(lambda x:
                                                                        '❌' if x < 0 else (
                                                                            '⚠️' if x < 92 else '✔️'))
    old_ba = ['Special Survey Remain Days', 'Intermediate Survey Remain Days', 'Docking Survey Remain Days', 'COC Remain Days']
    new_ba = ['Special', 'Intermediate', 'Docking', 'COC']
    for i,j in zip(new_ba, old_ba):
        df_diff_ba[i] = df_diff_ba[j].apply(lambda x:
                                            '❌' if x < 0 else (
                                            '⚠️' if x < 183 else '✔️'))

    data_mv = df_diff_mv[['Vessel Name', 'Class', 'Annual', 'Special', 'Intermediate', 'Docking', 'Tail Shaft', 'Boiler', 'COC',
                        'Annual Survey Remain Days2', 'Special Survey Remain Days2', 'Intermediate Survey Remain Days2', 'Docking Survey Remain Days2', 'Tail Shaft Survey Remain Days2', 'Boiler Survey Remain Days2', 'COC Remain Days2', 'Remarks']]
    data_cts = df_diff_cts[['Vessel Name', 'Class', 'Annual', 'Special', 'Intermediate', 'Docking', 'COC',
                            'Annual Survey Remain Days2', 'Special Survey Remain Days2', 'Intermediate Survey Remain Days2', 'Docking Survey Remain Days2', 'Cargo Gear Survey Remain Days2', 'COC Remain Days2', 'Remarks']]
    data_tb = df_diff_tb[['Vessel Name', 'Class', 'Annual', 'Special', 'Intermediate', 'Docking', 'Tail Shaft', 'COC',
                        'Annual Survey Remain Days2', 'Special Survey Remain Days2', 'Intermediate Survey Remain Days2', 'Docking Survey Remain Days2', 'Tail Shaft Survey Remain Days2', 'COC Remain Days2', 'Remarks']]
    data_ba = df_diff_ba[['Vessel Name', 'Class', 'Annual', 'Special', 'Intermediate', 'Docking', 'COC',
                        'Annual Survey Remain Days2', 'Special Survey Remain Days2', 'Intermediate Survey Remain Days2', 'Docking Survey Remain Days2', 'COC Remain Days2', 'Remarks']]
    data = pd.concat([data_mv, data_cts, data_tb, data_ba], keys=['OGV', 'CTS', 'Tugboat', 'Barge'])
    data.fillna({'Annual':'🔲', 'Special':'🔲', 'Intermediate':'🔲', 'Docking':'🔲', 'Tail Shaft':'🔲', 'Boiler':'🔲', 'COC':'🔲',
                'Annual Survey Remain Days2':'', 	
                'Special Survey Remain Days2':'',	
                'Intermediate Survey Remain Days2':'',	
                'Docking Survey Remain Days2':'',	
                'Tail Shaft Survey Remain Days2':'',	
                'Boiler Survey Remain Days2':'',
                'COC Remain Days2':'',
                'Cargo Gear Survey Remain Days2':'',
                'Remarks':''}, inplace=True)
    data = data.reset_index()

    ##### Table Image
    data_table = data[['level_0','Vessel Name','Class','Annual','Special','Intermediate','Docking','Tail Shaft','Boiler','COC','Remarks']]
    
    data_table = data_table.rename(columns={'level_0':'Type'})
    ogv_table = data_table[data_table['Type'] == 'OGV']
    cts_table = data_table[data_table['Type'] == 'CTS']
    tugboat_table = data_table[data_table['Type'] == 'Tugboat']
    barge_table = data_table[data_table['Type'] == 'Barge']

    ##### Header
    header_img = table_header()
    pio.write_image(header_img, 'img/header.png')

    ##### OGV
    # Mapping Color for each row and columns
    map_color = {"✔️":"#dfd", "🔲":"#dfd", "⚠️":"#ffd", "❌":"#fdd"}

    output = ['Annual','Special','Intermediate','Docking','Tail Shaft','Boiler','COC']
    output_col = ['Annual_col','Special_col','Intermediate_col','Docking_col','Tail Shaft_col','Boiler_col','COC_col']
    
    for i,j in zip(output, output_col):
        ogv_table[j] = ogv_table[i].map(map_color)
    
    # Insert color variable to a list (fill_color)
    fill_color_ogv = fill_color(ogv_table)

    # Table and save as img
    ogv_img = table_ogv(ogv_table, fill_color_ogv)
    pio.write_image(ogv_img, 'img/ogv.png')

    ##### CTS
    # Mapping Color for each row and columns
    for i,j in zip(output, output_col):
        cts_table[j] = cts_table[i].map(map_color)

    # Insert color variable to a list (fill_color)
    fill_color_cts = fill_color(cts_table)
            
    # Table and save as img
    cts_img = table_cts(cts_table, fill_color_cts)
    pio.write_image(cts_img, 'img/cts.png')

    ##### Tugboat
    # Mapping Color for each row and columns
    for i,j in zip(output, output_col):
        tugboat_table[j] = tugboat_table[i].map(map_color)

    # Insert color variable to a list (fill_color)
    fill_color_tugboat = fill_color(tugboat_table)
            
    # Table and save as img
    tugboat_img = table_tugboat(tugboat_table, fill_color_tugboat)
    pio.write_image(tugboat_img, 'img/tugboat.png')

    ##### Barge
    # Mapping Color for each row and columns
    for i,j in zip(output, output_col):
        barge_table[j] = barge_table[i].map(map_color)

    # Insert color variable to a list (fill_color)
    fill_color_barge = fill_color(barge_table)
            
    # Table and save as img
    barge_img = table_barge(barge_table, fill_color_barge)
    pio.write_image(barge_img, 'img/barge.png')

    #####
    path = os.getcwd()
    download_date = html.Strong('Downloaded slides as per ' + str(date.today().strftime("%d %b %Y")))
    ###################### PPTX ######################
    def to_pptx(bytes_io):
        # Create presentation
        pptx = path + '//' + 'slide_master.pptx'
        prs = Presentation(pptx)

        # define slidelayouts 
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide1 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide3 = prs.slides.add_slide(prs.slide_layouts[1])

        # title slide
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Certificate Monitoring Dashboard"
        subtitle.text = str(date.today().strftime("%d %b %Y"))

        # slide 1 OGV & CTS
        slide1.placeholders[10].text = 'ABL FLEET CLASS STATUS'
        slide1.shapes.add_picture('img/header.png', Cm(1.5), Cm(3.7), Cm(30.5), Cm(0.6))
        slide1.shapes.add_picture('img/ogv.png', Cm(1.5), Cm(4.3), Cm(30.5), Cm(2.4))

        slide1.shapes.add_picture('img/header.png', Cm(1.5), Cm(7.7), Cm(30.5), Cm(0.6))
        slide1.shapes.add_picture('img/cts.png', Cm(1.5), Cm(8.3), Cm(30.5), Cm(6.6))

        # slide 2 Tugboat
        slide2.placeholders[10].text = 'ABL FLEET CLASS STATUS'
        slide2.shapes.add_picture('img/header.png', Cm(1.5), Cm(3.7), Cm(30.5), Cm(0.6))
        slide2.shapes.add_picture('img/tugboat.png', Cm(1.5), Cm(4.3), Cm(30.5), Cm(12))

        # slide 3 Barge
        slide3.placeholders[10].text = 'ABL FLEET CLASS STATUS'
        slide3.shapes.add_picture('img/header.png', Cm(1.5), Cm(3.7), Cm(30.5), Cm(0.6))
        slide3.shapes.add_picture('img/barge.png', Cm(1.5), Cm(4.3), Cm(30.5), Cm(10.2))

        # Saving the PowerPoint presentation
        prs.save(bytes_io)    

        # return dcc.send_bytes(to_pptx, 'Equipment Dashboard.pptx')
    
    return dcc.send_bytes(to_pptx, 'Certificate Dashboard.pptx'), download_date


######################
# Function
######################
# 1. Function fill_color
def fill_color(table) :
# Insert color variable to a list (fill_color)
    fill = []
    output = ['Annual','Special','Intermediate','Docking','Tail Shaft','Boiler','COC']
    cols_to_show = ['Type','Vessel Name','Class','Annual','Special','Intermediate','Docking','Tail Shaft','Boiler','COC','Remarks']
    n = len(table)
    for col in cols_to_show:
        if col not in output:
            fill.append(['white']*n)
        else:
            fill.append(table[str(col)+'_col'].to_list())
    return(fill)


def table_header():
    fig = go.Figure(
        data=[go.Table(
            header=dict(values=[f"<b>{col}</b>" for col in ['VESSEL', 'CLASS SURVEY']],
                        fill_color='#2a3f5f',
                        line_color='white',
                        align='center',
                        font=dict(color='white', family="Verdana", size=20),
                        height=35,
                        ),
            columnwidth = [5,13],
        )
            ])
    fig.update_layout({'height':40,'width':2200,'plot_bgcolor': '#ffffff', 'paper_bgcolor': '#ffffff',
                        'margin' : {'t':0, 'b':0, 'l':0, 'r':0}, "autosize": True})
    return(fig)

# 2. Table OGV
def table_ogv(table, fill) :       
    # Table and save as img
    cols_to_show = ['Type','Vessel Name','Class','Annual','Special','Intermediate','Docking','Tail Shaft','Boiler','COC','Remarks']
    fig = go.Figure(
        data=[go.Table(
            header=dict(values=[f"<b>{col}</b>" for col in [x.upper() for x in cols_to_show]],
                        fill_color='#2a3f5f',
                        line_color='white',
                        align='center',
                        font=dict(color='white', family="Verdana", size=20),
                        height=35,
                        ),
            cells=dict(values=table[cols_to_show].values.T,
                    fill_color=fill,
                    line_color='lightgrey',
                    align='center',
                    font=dict(color='#2a3f5f', family="Verdana", size=20),
                    height=35
                    ),
            columnwidth = [1,3,1,1.5,1.5,1.5,1.5,1.5,1.5,1.5,2.5],
        )
            ])
    fig.update_layout({'height':150,'width':2200,'plot_bgcolor': '#ffffff', 'paper_bgcolor': '#ffffff',
                        'margin' : {'t':0, 'b':0, 'l':0, 'r':0}, "autosize": True})
    return(fig)

# 2. Table CTS
def table_cts(table, fill) : 
    # Table and save as img
    cols_to_show = ['Type','Vessel Name','Class','Annual','Special','Intermediate','Docking','Tail Shaft','Boiler','COC','Remarks']
    fig = go.Figure(
        data=[go.Table(
            header=dict(values=[f"<b>{col}</b>" for col in [x.upper() for x in cols_to_show]],
                        fill_color='#2a3f5f',
                        line_color='white',
                        align='center',
                        font=dict(color='white', family="Verdana", size=20),
                        height=35,
                        ),
            cells=dict(values=table[cols_to_show].values.T,
                    fill_color=fill,
                    line_color='lightgrey',
                    align='center',
                    font=dict(color='#2a3f5f', family="Verdana", size=20),
                    height=35
                    ),
            columnwidth = [1,3,1,1.5,1.5,1.5,1.5,1.5,1.5,1.5,2.5],
        )
            ])
    fig.update_layout({'height':390,'width':2200,'plot_bgcolor': '#ffffff', 'paper_bgcolor': '#ffffff',
                        'margin' : {'t':0, 'b':0, 'l':0, 'r':0}, "autosize": True})
    return(fig)

# 3. Table Tugboat
def table_tugboat(table, fill) : 
    # Table and save as img
    cols_to_show = ['Type','Vessel Name','Class','Annual','Special','Intermediate','Docking','Tail Shaft','Boiler','COC','Remarks']
    fig = go.Figure(
        data=[go.Table(
            header=dict(values=[f"<b>{col}</b>" for col in [x.upper() for x in cols_to_show]],
                        fill_color='#2a3f5f',
                        line_color='white',
                        align='center',
                        font=dict(color='white', family="Verdana", size=20),
                        height=35,
                        ),
            cells=dict(values=table[cols_to_show].values.T,
                    fill_color=fill,
                    line_color='lightgrey',
                    align='center',
                    font=dict(color='#2a3f5f', family="Verdana", size=20),
                    height=35
                    ),
            columnwidth = [1,3,1,1.5,1.5,1.5,1.5,1.5,1.5,1.5,2.5],
        )
            ])
    fig.update_layout({'height':720,'width':2200,'plot_bgcolor': '#ffffff', 'paper_bgcolor': '#ffffff',
                        'margin' : {'t':0, 'b':0, 'l':0, 'r':0}, "autosize": True})
    return(fig)

# 4. Table Barge
def table_barge(table, fill) : 
    # Table and save as img
    cols_to_show = ['Type','Vessel Name','Class','Annual','Special','Intermediate','Docking','Tail Shaft','Boiler','COC','Remarks']
    fig = go.Figure(
        data=[go.Table(
            header=dict(values=[f"<b>{col}</b>" for col in [x.upper() for x in cols_to_show]],
                        fill_color='#2a3f5f',
                        line_color='white',
                        align='center',
                        font=dict(color='white', family="Verdana", size=20),
                        height=35,
                        ),
            cells=dict(values=table[cols_to_show].values.T,
                    fill_color=fill,
                    line_color='lightgrey',
                    align='center',
                    font=dict(color='#2a3f5f', family="Verdana", size=20),
                    height=35
                    ),
            columnwidth = [1,3,1,1.5,1.5,1.5,1.5,1.5,1.5,1.5,2.5],
        )
            ])
    fig.update_layout({'height':610,'width':2200,'plot_bgcolor': '#ffffff', 'paper_bgcolor': '#ffffff',
                        'margin' : {'t':0, 'b':0, 'l':0, 'r':0}, "autosize": True})
    return(fig)